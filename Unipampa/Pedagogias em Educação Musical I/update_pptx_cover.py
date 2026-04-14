import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Paths
base_dir = "/Users/rafaelrodriguesdasilva/Documents/Agentes - Antigravity/Pedagogias em Educação Musical I"
pptx_path = os.path.join(base_dir, "Apresentacao_Dalcroze.pptx")
output_path = os.path.join(base_dir, "Apresentacao_Dalcroze_Final.pptx")
photo_path = "/Users/rafaelrodriguesdasilva/Documents/Agentes - Antigravity/rafaeldasilva-pro/public/assets/uploaded_image_2_1765859024622.jpg"
logo_path = os.path.join(base_dir, "unipampa_logo.png")
synthesis_path = "/Users/rafaelrodriguesdasilva/.gemini/antigravity/brain/e6bb69a9-875b-4a08-bb30-46a8a80ee785/dalcroze_visual_synthesis_1773705540552.png"

def update_presentation():
    # Load presentation
    if not os.path.exists(pptx_path):
        print(f"Error: {pptx_path} not found.")
        return

    prs = Presentation(pptx_path)

    # Add a new blank slide at the beginning
    blank_slide_layout = prs.slide_layouts[6]  # Usually blank
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Move the new slide to the first position
    # python-pptx doesn't have a direct move_slide, so we'll just insert everything on index 0
    # Actually, prs.slides.add_slide adds to the end. 
    # Let's try to modify slide 0 instead.
    if len(prs.slides) > 0:
        slide = prs.slides[0]
        # Clear existing shapes on the first slide
        for shape in list(slide.shapes):
            sp = shape._element
            sp.getparent().remove(sp)
    else:
        slide = prs.slides.add_slide(blank_slide_layout)

    # Slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # 1. Background Visual Synthesis (Full Slide)
    if os.path.exists(synthesis_path):
        slide.shapes.add_picture(synthesis_path, 0, 0, width=slide_width, height=slide_height)

    # 2. Professor Rafael's Photo (Lower Right Corner)
    if os.path.exists(photo_path):
        photo_width = Inches(2.0)
        photo_height = Inches(2.0)
        left = slide_width - photo_width - Inches(0.5)
        top = slide_height - photo_height - Inches(0.5)
        slide.shapes.add_picture(photo_path, left, top, width=photo_width, height=photo_height)

    # 3. Unipampa Logo (Lower Left Corner)
    if os.path.exists(logo_path):
        logo_width = Inches(1.8)
        left = Inches(0.5)
        top = slide_height - Inches(1.2)
        slide.shapes.add_picture(logo_path, left, top, width=logo_width)

    # 4. Professional Title
    title_box = slide.shapes.add_textbox(Inches(0.5), slide_height - Inches(2.0), slide_width - Inches(1.0), Inches(0.5))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "Prof. Dr. Rafael Rodrigues da Silva"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(255, 255, 255)  # White text for contrast

    # Save
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    update_presentation()
